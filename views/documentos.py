import streamlit as st
import pandas as pd
from PyPDF2 import PdfReader
import docx
from io import BytesIO
from docx import Document
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Inches
import tempfile
import plotly.express as px
import plotly.io as pio

# ---------- Funções de leitura ----------
def ler_txt(file):
    return file.read().decode("utf-8")

def ler_pdf(file):
    pdf = PdfReader(file)
    texto = ""
    for page in pdf.pages:
        page_text = page.extract_text()
        if page_text:
            texto += page_text + "\n"
    return texto

def ler_docx(file):
    doc = docx.Document(file)
    texto = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    return texto

def resumir_texto(texto, max_chars=500):
    if len(texto) <= max_chars:
        return texto
    return texto[:max_chars] + "..."

# ---------- Inicializa session_state ----------
if "document_links" not in st.session_state:
    st.session_state.document_links = []

# ---------- Função principal da página ----------
def show_documentos():
    # ---------- Estilo Chinês ----------
    st.markdown("""
    <style>
    body { background-color: #FFF8E7; font-family: 'Noto Serif SC', serif; }
    h1,h2,h3 { color: #B22222; text-shadow: 1px 1px 2px #FFD700; }
    .stButton>button { background-color: #FF6347; color:white; font-weight:bold; border-radius:10px; }
    </style>
    """, unsafe_allow_html=True)

    # ---------- Dados fictícios do dashboard ----------
    df_exemplo = pd.DataFrame({
        "Nome": ["Ana", "Bruno", "Carlos", "Diana"],
        "Média": [8.5, 6.2, 7.9, 5.4]
    })
    fig_bar = px.bar(df_exemplo, x="Nome", y="Média", text="Média",
                     color="Média", color_continuous_scale=px.colors.sequential.Oranges,
                     title="Média Individual dos Estudantes 🐉")
    fig_bar.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    fig_pie = px.pie(df_exemplo, names="Nome", values="Média", title="Distribuição das Médias 🏮",
                     color_discrete_sequence=["#B22222","#FF8C00","#FFD700"])

    # ---------- Interface ----------
    st.title("📊 Análise de Documentos 🏮🐉")
    st.write("Envie um documento (.txt, .pdf ou .docx) e veja uma descrição gerada automaticamente.")

    uploaded_file = st.file_uploader("📂 Envie seu documento", type=["txt", "pdf", "docx"])

    if uploaded_file:
        st.markdown(f"### 📑 {uploaded_file.name}")
        texto = ""
        file_type = uploaded_file.type

        try:
            if uploaded_file.name.endswith(".txt") or file_type == "text/plain":
                texto = ler_txt(uploaded_file)
            elif uploaded_file.name.endswith(".pdf") or file_type == "application/pdf":
                texto = ler_pdf(uploaded_file)
            elif uploaded_file.name.endswith(".docx") or file_type in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword"
            ]:
                texto = ler_docx(uploaded_file)
            else:
                st.error("Formato de arquivo não suportado.")
        except Exception as e:
            st.error(f"Erro ao ler o arquivo: {e}")

        if texto:
            st.subheader("📖 Pré-visualização do conteúdo")
            st.text_area("Conteúdo extraído", texto[:2000], height=200)

            st.subheader("📝 Descrição do documento")
            st.write(resumir_texto(texto))

            # ---------- Botões de geração ----------
            col1, col2, col3 = st.columns(3)

            with col1:
                if st.button("📑 Gerar DOCX"):
                    doc = Document()
                    doc.add_heading(uploaded_file.name, 0)
                    doc.add_paragraph(texto)
                    buffer = BytesIO()
                    doc.save(buffer)
                    buffer.seek(0)
                    st.session_state.document_links.append({
                        "nome": f"{uploaded_file.name}.docx",
                        "buffer": buffer,
                        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    })
                    st.success("DOCX gerado!")

            with col2:
                if st.button("📄 Gerar PDF com gráficos"):
                    buffer = BytesIO()
                    doc_pdf = SimpleDocTemplate(buffer, pagesize=letter)
                    elements = []
                    styles = getSampleStyleSheet()
                    elements.append(Paragraph(uploaded_file.name, styles['Title']))
                    elements.append(Spacer(1,12))
                    elements.append(Paragraph(texto, styles['Normal']))
                    elements.append(Spacer(1,12))

                    # Adiciona gráficos
                    for fig, title in [(fig_bar, "Média Individual dos Estudantes 🐉"), 
                                       (fig_pie, "Distribuição das Médias 🏮")]:
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
                            pio.write_image(fig, tmpfile.name, format='png', scale=2)
                            elements.append(Paragraph(title, styles['Heading2']))
                            elements.append(Image(tmpfile.name, width=400, height=300))
                            elements.append(Spacer(1,12))

                    doc_pdf.build(elements)
                    buffer.seek(0)
                    st.session_state.document_links.append({
                        "nome": f"{uploaded_file.name}_com_graficos.pdf",
                        "buffer": buffer,
                        "mime": "application/pdf"
                    })
                    st.success("PDF com gráficos gerado!")

            with col3:
                if st.button("📊 Gerar PowerPoint com gráficos"):
                    prs = Presentation()
                    slide_layout = prs.slide_layouts[5]

                    # Slide texto
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = uploaded_file.name
                    slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4)).text = texto

                    # Slides de gráficos
                    for fig, title in [(fig_bar, "Média Individual dos Estudantes 🐉"), 
                                       (fig_pie, "Distribuição das Médias 🏮")]:
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
                            pio.write_image(fig, tmpfile.name, format='png', scale=2)
                            slide = prs.slides.add_slide(slide_layout)
                            slide.shapes.title.text = title
                            slide.shapes.add_picture(tmpfile.name, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

                    buffer = BytesIO()
                    prs.save(buffer)
                    buffer.seek(0)
                    st.session_state.document_links.append({
                        "nome": f"{uploaded_file.name}_com_graficos.pptx",
                        "buffer": buffer,
                        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    })
                    st.success("PowerPoint com gráficos gerado!")

            # ---------- Exibe links gerados ----------
            if st.session_state.document_links:
                st.subheader("📂 Documentos gerados")
                for doc_info in st.session_state.document_links:
                    st.download_button(
                        label=f"📥 Baixar {doc_info['nome']}",
                        data=doc_info["buffer"],
                        file_name=doc_info["nome"],
                        mime=doc_info["mime"]
                    )
