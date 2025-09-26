import streamlit as st
import pandas as pd
import plotly.express as px
from controllers.estudante_controller import listar_estudantes

# Imports para exportação
from io import BytesIO
from docx import Document
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
import plotly.io as pio
import tempfile
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(
    page_title="📊 Dashboard Escolar 🐉🏮",
    page_icon="🐲",
    layout="wide"
)

# ---------- HTML + CSS Interativo ----------
st.markdown("""
<style>
body { font-family: 'Noto Serif SC', serif; background-color: #FFF8E7; overflow-x: hidden;}
.stButton>button { background-color: #FF6347; color:white; border-radius:10px; font-weight:bold;}
.stTable td, .stTable th { border-color:#B22222 !important; }
</style>
""", unsafe_allow_html=True)

# ---------- Dashboard ----------
def show():
    st.title("📊 Dashboard Escolar 🐉🏮")
    st.write("Relatórios e estatísticas dos alunos aqui...")

    listaDeEstudantes = listar_estudantes()
    if not listaDeEstudantes:
        st.warning("Nenhum estudante foi cadastrado ainda")
        return

    df = pd.DataFrame([{"Nome": e.nome,"1º Nota": e.nota1,"2º Nota": e.nota2,"Média": e.media} for e in listaDeEstudantes])

    # ---------- FILTROS ----------
    st.sidebar.header("Filtros do Dashboard 🏮")
    faixa_media = st.sidebar.slider("Selecione a faixa de média", 0.0, 10.0, (0.0,10.0), 0.1)
    nome_filtrado = st.sidebar.text_input("Buscar por nome do estudante 🐉")

    df_filtrado = df[(df["Média"] >= faixa_media[0]) & (df["Média"] <= faixa_media[1])]
    if nome_filtrado:
        df_filtrado = df_filtrado[df_filtrado["Nome"].str.contains(nome_filtrado, case=False)]

    # ---------- Tabela ----------
    st.subheader("Lista de Estudantes Cadastrados 🏮")
    st.table(df_filtrado.style.format({"1º Nota":"{:.2f}","2º Nota":"{:.2f}","Média":"{:.2f}"}))

    # ---------- Gráficos ----------
    fig_bar = px.bar(df_filtrado, x="Nome", y="Média", title="Média Individual dos Estudantes 🐉",
                     labels={"Média":"Média","Nome":"Estudante"}, text="Média",
                     range_y=[0,10], color="Média", color_continuous_scale=px.colors.sequential.Oranges)
    fig_bar.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    st.plotly_chart(fig_bar, use_container_width=True)

    categorias = {
        "Abaixo de 5": len(df_filtrado[df_filtrado["Média"] < 5]),
        "Entre 5 e 7": len(df_filtrado[(df_filtrado["Média"] >= 5) & (df_filtrado["Média"] <= 7)]),
        "Acima de 7": len(df_filtrado[df_filtrado["Média"] > 7])
    }
    fig_pie = px.pie(names=list(categorias.keys()), values=list(categorias.values()),
                     title="Distribuição das Médias dos Estudantes 🏮",
                     color_discrete_sequence=["#B22222","#FF8C00","#FFD700"])
    st.plotly_chart(fig_pie, use_container_width=True)

    # ---------- Exportar Dashboard ----------
    st.subheader("📂 Exportar Dashboard")

    # DOCX
    if st.button("📑 Exportar para Word (DOCX)"):
        doc = Document()
        doc.add_heading("Relatório Escolar - Dashboard 🐉🏮", 0)
        doc.add_paragraph("Relatório gerado a partir do sistema interativo.")

        t = doc.add_table(rows=1, cols=len(df_filtrado.columns))
        hdr_cells = t.rows[0].cells
        for i, col in enumerate(df_filtrado.columns):
            hdr_cells[i].text = col
        for _, row in df_filtrado.iterrows():
            row_cells = t.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = str(val)

        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        st.download_button(
            label="📥 Baixar Relatório Word",
            data=buffer,
            file_name="dashboard_escolar.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    # PDF
    if st.button("📄 Exportar para PDF"):
        buffer = BytesIO()
        doc_pdf = SimpleDocTemplate(buffer, pagesize=letter)
        elements = []
        styles = getSampleStyleSheet()

        elements.append(Paragraph("Relatório Escolar - Dashboard 🐉🏮", styles['Title']))
        elements.append(Paragraph("Relatório gerado a partir do sistema interativo.", styles['Normal']))
        elements.append(Spacer(1,12))

        data_table = [list(df_filtrado.columns)] + df_filtrado.values.tolist()
        t = Table(data_table)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#B22222")),
            ('TEXTCOLOR',(0,0),(-1,0),colors.white),
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('GRID', (0,0), (-1,-1), 1, colors.black),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold')
        ]))
        elements.append(t)
        elements.append(Spacer(1,12))

        for fig, title in [(fig_bar, "Média Individual dos Estudantes 🐉"), (fig_pie, "Distribuição das Médias dos Estudantes 🏮")]:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
                pio.write_image(fig, tmpfile.name, format='png', scale=2)
                elements.append(Paragraph(title, styles['Heading2']))
                elements.append(Image(tmpfile.name, width=400, height=300))
                elements.append(Spacer(1,12))

        doc_pdf.build(elements)
        buffer.seek(0)

        st.download_button(
            label="📥 Baixar Relatório PDF Completo",
            data=buffer,
            file_name="dashboard_escolar_completo.pdf",
            mime="application/pdf"
        )

    # PowerPoint
    if st.button("📊 Exportar para PowerPoint (PPTX)"):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # layout em branco

        # Slide Tabela
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Tabela de Estudantes"
        rows, cols = df_filtrado.shape
        top, left, width, height = Inches(1), Inches(0.5), Inches(9), Inches(5)
        table = slide.shapes.add_table(rows=rows+1, cols=cols, left=left, top=top, width=width, height=height).table

        # Cabeçalho
        for j, col_name in enumerate(df_filtrado.columns):
            table.cell(0, j).text = str(col_name)

        # Dados
        for i in range(rows):
            for j in range(cols):
                table.cell(i+1, j).text = str(df_filtrado.iloc[i, j])

        # Slide Gráfico de barras
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
            pio.write_image(fig_bar, tmpfile.name, format='png', scale=2)
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Média Individual dos Estudantes 🐉"
            slide.shapes.add_picture(tmpfile.name, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

        # Slide Gráfico de pizza
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
            pio.write_image(fig_pie, tmpfile.name, format='png', scale=2)
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Distribuição das Médias dos Estudantes 🏮"
            slide.shapes.add_picture(tmpfile.name, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        st.download_button(
            label="📥 Baixar Dashboard em PowerPoint",
            data=buffer,
            file_name="dashboard_escolar.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

